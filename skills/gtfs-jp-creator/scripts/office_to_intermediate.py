#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
office_to_intermediate.py
=========================
Step1 前段(Office経路): Word(.docx) / PowerPoint(.pptx) の時刻表を、
既存の抽出経路にそのまま乗せられる中間ファイルへ変換する。

なぜ「変換」なのか（今までの生成データに合わせる）:
  - バス時刻表は「表(テーブル)」で作られていることが多い。表はセルがグリッド構造で
    機械可読なので、Excel経路(extract_timetable_excel.py)と同じく **直接セルを読む**のが
    最も確実・高精度（PDF座標クラスタリングやOCRが不要）。
  - よって docx/pptx の中の「時刻の入った表」を **.xlsx に写して** Excel経路へ渡す。
  - 表ではなく「時刻表の画像を貼り付けただけ」の場合は、画像を取り出して 1つの **PDF**に
    まとめ、既存の画像PDF→OCR(MinerU)経路へ渡す。

設計方針(正しく失敗する):
  - 機械的に取り出せるもの(表のセル文字・埋め込み画像)だけを扱う。
  - 便名・方向・循環などの解釈はしない(Step2に委ねる)。
  - 時刻らしい表も画像も無ければ kind="none" を返し、推測でデータを作らない。

出力(標準出力に JSON 1行):
  {"kind":"xlsx","paths":["out/office_1.xlsx", ...]}   # 時刻表の表を写した(表ごと1ファイル)
  {"kind":"pdf","path":"out/office_images.pdf"}         # 画像を束ねたPDF(→OCRへ)
  {"kind":"none","message":"..."}                       # 時刻表が見つからない

Usage:
  python office_to_intermediate.py <input.docx|.pptx> --outdir <dir>

License: Apache 2.0
"""
import argparse
import json
import re
import sys
from pathlib import Path

# HH:MM / H:MM / HH：MM(全角) を時刻とみなす（Excel経路と同じ発想の緩い判定）
_TIME_RE = re.compile(r"\b\d{1,2}\s*[:：]\s*\d{2}\b")


def _looks_timetable(grid):
    """2次元セル配列に時刻セルが複数あれば時刻表とみなす（誤検出を避け 2件以上）。"""
    n = 0
    for row in grid:
        for cell in row:
            if cell and _TIME_RE.search(str(cell)):
                n += 1
                if n >= 2:
                    return True
    return False


def _docx_tables(path):
    """docx の全テーブルを 2次元配列(list[list[str]])で返す。結合セルは原文のまま。"""
    import docx
    doc = docx.Document(str(path))
    grids = []
    for tbl in doc.tables:
        grid = []
        for row in tbl.rows:
            grid.append([(c.text or "").strip() for c in row.cells])
        if grid:
            grids.append(grid)
    return grids


def _docx_images(path, outdir):
    """docx 内の埋め込み画像を outdir に保存し、パス一覧を返す。"""
    import docx
    doc = docx.Document(str(path))
    saved = []
    for i, rel in enumerate(doc.part.rels.values()):
        if "image" in rel.reltype:
            try:
                blob = rel.target_part.blob
                ext = (rel.target_ref.rsplit(".", 1)[-1] or "png").lower()
                p = Path(outdir) / f"office_img_{i}.{ext}"
                p.write_bytes(blob)
                saved.append(p)
            except Exception:
                continue
    return saved


def _pptx_tables(path):
    """pptx の全スライドの表を 2次元配列で返す。結合の被セルは空文字にする。"""
    from pptx import Presentation
    prs = Presentation(str(path))
    grids = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            tbl = shape.table
            grid = []
            for r in range(len(tbl.rows)):
                row = []
                for c in range(len(tbl.columns)):
                    cell = tbl.cell(r, c)
                    # 結合セルの被セル(spanned)は空にして時刻の重複を防ぐ
                    if getattr(cell, "is_spanned", False):
                        row.append("")
                    else:
                        row.append((cell.text or "").strip())
                grid.append(row)
            if grid:
                grids.append(grid)
    return grids


def _pptx_images(path, outdir):
    """pptx 内の画像シェイプを outdir に保存し、パス一覧を返す。"""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(str(path))
    saved = []
    i = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    p = Path(outdir) / f"office_img_{i}.{img.ext or 'png'}"
                    p.write_bytes(img.blob)
                    saved.append(p)
                    i += 1
                except Exception:
                    continue
    return saved


def _write_xlsx(grid, out_path):
    """2次元配列を .xlsx に書き出す（Excel経路がそのまま読める形）。"""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "timetable"
    for r, row in enumerate(grid, start=1):
        for c, val in enumerate(row, start=1):
            if val:
                ws.cell(row=r, column=c, value=str(val))
    wb.save(str(out_path))


def _images_to_pdf(images, out_path):
    """画像群を1つのPDFにまとめる（画像PDF→OCR経路に渡すため）。"""
    from PIL import Image
    pages = []
    for ip in images:
        try:
            im = Image.open(str(ip)).convert("RGB")
            pages.append(im)
        except Exception:
            continue
    if not pages:
        return False
    pages[0].save(str(out_path), "PDF", save_all=True, append_images=pages[1:])
    return True


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("input")
    ap.add_argument("--outdir", required=True)
    a = ap.parse_args()
    src = Path(a.input)
    outdir = Path(a.outdir)
    outdir.mkdir(parents=True, exist_ok=True)
    low = src.name.lower()

    if low.endswith(".docx"):
        grids = _docx_tables(src)
        get_images = _docx_images
    elif low.endswith(".pptx"):
        grids = _pptx_tables(src)
        get_images = _pptx_images
    else:
        print(json.dumps({"kind": "none",
                          "message": "docx / pptx 以外は対象外です。"}, ensure_ascii=False))
        return

    # 1) 時刻の入った表 → 表ごとに .xlsx（複数路線は既存の複数ファイル統合に乗る）
    tt = [g for g in grids if _looks_timetable(g)]
    if tt:
        paths = []
        for i, g in enumerate(tt, start=1):
            p = outdir / f"office_{i}.xlsx"
            _write_xlsx(g, p)
            paths.append(str(p))
        print(json.dumps({"kind": "xlsx", "paths": paths}, ensure_ascii=False))
        return

    # 2) 表が無い → 貼り付け画像を1つのPDFにまとめOCR経路へ
    imgs = get_images(src, outdir)
    if imgs:
        pdf = outdir / "office_images.pdf"
        if _images_to_pdf(imgs, pdf):
            print(json.dumps({"kind": "pdf", "path": str(pdf)}, ensure_ascii=False))
            return

    # 3) どちらも無い → 正しく失敗
    print(json.dumps({"kind": "none",
                      "message": "時刻表の表も画像も見つかりませんでした。"
                                 "Word/PowerPoint内に時刻表を『表』または『画像』として入れてください。"},
                     ensure_ascii=False))


if __name__ == "__main__":
    main()
